#!/usr/bin/env python3
"""
Generate a 5-minute undergraduate thesis defense deck.
"""
from pathlib import Path

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parents[2]
THESIS_DIR = ROOT / "docs" / "thesis"
FIG_DIR = THESIS_DIR / "figures"
OUT_DIR = THESIS_DIR / "defense"
OUT_DIR.mkdir(parents=True, exist_ok=True)

PPTX_PATH = OUT_DIR / "本科毕设答辩_5分钟_KAN超图药物协同预测.pptx"
NOTES_PATH = OUT_DIR / "本科毕设答辩_5分钟_讲稿.md"
METRIC_FIG = OUT_DIR / "cold_drug_metrics.png"


WIDE_W = Cm(33.867)
WIDE_H = Cm(19.05)

COLORS = {
    "ink": RGBColor(25, 34, 45),
    "muted": RGBColor(92, 105, 120),
    "line": RGBColor(218, 224, 231),
    "bg": RGBColor(248, 250, 252),
    "panel": RGBColor(255, 255, 255),
    "navy": RGBColor(24, 62, 98),
    "teal": RGBColor(16, 133, 122),
    "orange": RGBColor(215, 117, 33),
    "red": RGBColor(184, 73, 72),
    "green": RGBColor(47, 132, 94),
}


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def set_line(shape, color, width=1):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_textbox(slide, text, left, top, width, height, font_size=20, color=None,
                bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
                font_name="Microsoft YaHei"):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Cm(0.05)
    frame.margin_right = Cm(0.05)
    frame.margin_top = Cm(0.02)
    frame.margin_bottom = Cm(0.02)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]
    return box


def add_multiline(slide, lines, left, top, width, height, font_size=18,
                  color=None, bullet=False, line_spacing=1.08):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.line_spacing = line_spacing
        if bullet:
            p.text = "• " + line
        for run in p.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color or COLORS["ink"]
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, title, Cm(1.3), Cm(0.7), Cm(26), Cm(1.05),
                font_size=28, color=COLORS["ink"], bold=True)
    if subtitle:
        add_textbox(slide, subtitle, Cm(1.35), Cm(1.75), Cm(25), Cm(0.55),
                    font_size=13, color=COLORS["muted"])
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(1.35), Cm(2.45), Cm(31.1), Cm(0.03)
    )
    set_fill(line, COLORS["line"])


def add_page_no(slide, no):
    add_textbox(slide, f"{no:02d}", Cm(31.1), Cm(17.8), Cm(1.1), Cm(0.45),
                font_size=10, color=COLORS["muted"], align=PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, title=None, accent=None):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    card.adjustments[0] = 0.08
    set_fill(card, COLORS["panel"])
    set_line(card, COLORS["line"], 0.8)
    if accent:
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                     left, top, Cm(0.12), height)
        set_fill(bar, accent)
    if title:
        add_textbox(slide, title, left + Cm(0.45), top + Cm(0.28),
                    width - Cm(0.8), Cm(0.55), font_size=15,
                    color=accent or COLORS["navy"], bold=True)
    return card


def add_metric(slide, value, label, left, top, width, accent):
    add_card(slide, left, top, width, Cm(2.2), accent=accent)
    add_textbox(slide, value, left + Cm(0.4), top + Cm(0.35),
                width - Cm(0.8), Cm(0.75), font_size=24, color=accent,
                bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, left + Cm(0.25), top + Cm(1.25),
                width - Cm(0.5), Cm(0.55), font_size=12,
                color=COLORS["muted"], align=PP_ALIGN.CENTER)


def add_picture_fit(slide, image_path, left, top, width, height):
    from PIL import Image
    image_path = Path(image_path)
    img = Image.open(image_path)
    img_w, img_h = img.size
    box_ratio = width / height
    img_ratio = img_w / img_h
    if img_ratio > box_ratio:
        pic_w = width
        pic_h = width / img_ratio
        pic_left = left
        pic_top = top + (height - pic_h) / 2
    else:
        pic_h = height
        pic_w = height * img_ratio
        pic_left = left + (width - pic_w) / 2
        pic_top = top
    return slide.shapes.add_picture(str(image_path), pic_left, pic_top,
                                    width=pic_w, height=pic_h)


def add_table(slide, rows, cols, left, top, width, height, data,
              col_widths=None, font_size=11):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            cell.margin_left = Cm(0.08)
            cell.margin_right = Cm(0.08)
            cell.margin_top = Cm(0.05)
            cell.margin_bottom = Cm(0.05)
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = COLORS["navy"] if r == 0 else RGBColor(255, 255, 255)
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = "Microsoft YaHei"
                run.font.size = Pt(font_size if r else font_size + 1)
                run.font.bold = r == 0
                run.font.color.rgb = RGBColor(255, 255, 255) if r == 0 else COLORS["ink"]
    return table_shape


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["bg"]
    return slide


def make_metric_chart():
    models = ["MLP baseline", "KAN head", "DrugKAN", "HgKAN-Agg"]
    aupr = [0.2758, 0.1717, 0.1967, 0.2765]
    auc = [0.7333, 0.5590, 0.6485, 0.7452]
    rmse = [18.3997, 19.2865, 17.8930, 17.6763]
    mae = [13.6750, 14.1486, 13.1844, 12.8615]

    fig, axes = plt.subplots(1, 2, figsize=(10.6, 3.5), dpi=180)
    bar_colors = ["#6c7a89", "#b84948", "#d77521", "#10857a"]
    x = range(len(models))

    axes[0].bar([i - 0.18 for i in x], aupr, width=0.34, label="AUPR", color=bar_colors)
    axes[0].bar([i + 0.18 for i in x], auc, width=0.34, label="AUC", color="#8fb8de")
    axes[0].set_title("Classification: higher is better", fontsize=11, weight="bold")
    axes[0].set_ylim(0, 0.82)
    axes[0].set_xticks(list(x), models, rotation=16, ha="right", fontsize=8)
    axes[0].grid(axis="y", alpha=0.22)
    axes[0].legend(frameon=False, fontsize=8)

    axes[1].bar([i - 0.18 for i in x], rmse, width=0.34, label="RMSE", color=bar_colors)
    axes[1].bar([i + 0.18 for i in x], mae, width=0.34, label="MAE", color="#8fb8de")
    axes[1].set_title("Regression: lower is better", fontsize=11, weight="bold")
    axes[1].set_ylim(12, 20.4)
    axes[1].set_xticks(list(x), models, rotation=16, ha="right", fontsize=8)
    axes[1].grid(axis="y", alpha=0.22)
    axes[1].legend(frameon=False, fontsize=8)

    for ax in axes:
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.spines["left"].set_color("#d8dee6")
        ax.spines["bottom"].set_color("#d8dee6")

    fig.suptitle("Cold-drug split main results", fontsize=13, weight="bold", y=1.04)
    fig.tight_layout()
    fig.savefig(METRIC_FIG, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def add_notes():
    notes = """# 5分钟答辩讲稿

## 1. 封面（约20秒）
各位老师好，我的题目是“基于KAN超图的药物协同预测研究”。本文关注的问题是：在药物组合筛选中，如何利用药物分子结构、细胞系表达和两药组合关系，预测某一组合是否具有协同效应。

## 2. 研究问题（约40秒）
联合用药空间会随着药物数量快速膨胀，完全依赖湿实验成本很高。这个任务难在三点：第一，正例比例低，本文数据中协同比例约14.8%；第二，random划分容易高估模型能力；第三，真实应用更关心未见药物的cold-drug泛化。因此本文不只问模型能不能拟合历史样本，而是问KAN应该放在哪里才可能对新药组合筛选有帮助。

## 3. 数据与任务（约40秒）
本文构建约16万条DrugComb相关样本，输入包括两个药物分子图和一个细胞系表达向量，输出包括协同二分类标签和连续Loewe分数。评价上，分类主要看AUPR和AUC，回归看RMSE和MAE。cold-drug划分中训练集和测试集药物集合重叠为0，避免药物实体泄漏。

## 4. 方法（约60秒）
整体框架由三部分组成：药物图编码器、细胞表达编码器和drug-drug-cell交互模块。本文比较了多种KAN放置方式：直接替换最终MLP的KAN head；放在药物图编码器中的DrugKAN；以及放在结构化交互层中的HgKAN。最终最有效的是HgKAN-Agg：先把两种药物聚合成drug-pair节点，再与cell节点通过KAN消息函数交互。

## 5. 实验结果（约70秒）
random split下MLP baseline仍然很强，说明KAN不是通用替代品。关键结果在cold-drug：直接KAN head明显退化，分类AUC只有0.5590；DrugKAN在回归中有一定收益，但分类不稳定；HgKAN-Agg在分类AUPR/AUC和回归RMSE/MAE上均优于MLP baseline。尤其回归RMSE从18.3997降到17.6763，说明药物对聚合图给KAN提供了更合适的结构载体。

## 6. 可解释性案例（约50秒）
案例是Lenalidomide和mitomycin C在IGROV1细胞系上的cold-drug正样本，两种药都未出现在训练集中。模型预测协同概率为0.9106。屏蔽两种药物后概率降到0.3711；当药物对表示从0恢复到15%时，预测概率越过0.5阈值。这说明模型确实依赖药物对输入。原子saliency还显示高贡献区域与Lenalidomide的glutarimide环、mitomycin C的quinone/mitosene骨架有合理重合。

## 7. 结论（约40秒）
本文的结论不是“KAN全面优于MLP”，而是：KAN的效果高度依赖放置位置。无结构地替换最终预测头效果较差；把KAN作为图或超图消息传递中的可学习非线性函数更合理。本文的贡献在于系统比较KAN放置方式，得到HgKAN-Agg在cold-drug场景下的正结果，并构建了函数曲线、扰动和原子saliency组成的解释流程。
"""
    NOTES_PATH.write_text(notes, encoding="utf-8")


def build_deck():
    make_metric_chart()

    prs = Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H

    # Slide 1
    slide = blank_slide(prs)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(33.867), Cm(19.05))
    set_fill(band, RGBColor(239, 246, 247))
    side = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(0.42), Cm(19.05))
    set_fill(side, COLORS["teal"])
    add_textbox(slide, "基于KAN超图的药物协同预测研究",
                Cm(1.7), Cm(4.0), Cm(24.5), Cm(1.35), font_size=32,
                color=COLORS["ink"], bold=True)
    add_textbox(slide, "本科毕业论文答辩 · 5分钟版",
                Cm(1.75), Cm(5.55), Cm(13.5), Cm(0.6), font_size=16,
                color=COLORS["teal"], bold=True)
    add_textbox(slide, "答辩人：XXX    指导教师：XXX    专业：XXX",
                Cm(1.75), Cm(15.3), Cm(18), Cm(0.55), font_size=14,
                color=COLORS["muted"])
    add_metric(slide, "16万+", "药物组合样本", Cm(22.1), Cm(4.2), Cm(4.4), COLORS["teal"])
    add_metric(slide, "0泄漏", "cold-drug药物重叠", Cm(27.0), Cm(4.2), Cm(4.4), COLORS["orange"])
    add_metric(slide, "HgKAN-Agg", "主结果模型", Cm(22.1), Cm(7.2), Cm(9.3), COLORS["navy"])
    add_page_no(slide, 1)

    # Slide 2
    slide = blank_slide(prs)
    add_title(slide, "研究问题：药物协同是三元交互预测", "目标不是替代湿实验，而是为候选组合排序与优先级筛选提供模型依据")
    add_card(slide, Cm(1.35), Cm(3.0), Cm(9.5), Cm(11.8), "任务挑战", COLORS["red"])
    add_multiline(slide, [
        "组合空间随药物数量呈二次增长",
        "协同正例稀少，本文正例比例约14.8%",
        "random split容易利用已见药物先验",
        "真实筛选更关注未见药物的cold-drug泛化",
        "药物发现场景需要可解释线索"
    ], Cm(1.85), Cm(4.15), Cm(8.5), Cm(8.5), font_size=16, bullet=True)
    add_card(slide, Cm(12.0), Cm(3.0), Cm(9.5), Cm(11.8), "核心问题", COLORS["teal"])
    add_multiline(slide, [
        "KAN能否提升药物协同预测？",
        "KAN应放在最终预测头、药物图编码器，还是drug-drug-cell交互层？",
        "模型是否真正依赖药物对结构，而非只学习数据先验？"
    ], Cm(12.5), Cm(4.15), Cm(8.4), Cm(8.5), font_size=17, bullet=True)
    add_card(slide, Cm(22.65), Cm(3.0), Cm(9.8), Cm(11.8), "本文答案", COLORS["navy"])
    add_textbox(slide, "KAN不是MLP的通用替代品",
                Cm(23.2), Cm(4.25), Cm(8.6), Cm(0.8), font_size=18,
                color=COLORS["red"], bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, "更适合作为结构化消息传递中的可学习非线性函数",
                Cm(23.2), Cm(6.1), Cm(8.6), Cm(1.6), font_size=21,
                color=COLORS["navy"], bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, "结构先行，KAN局部函数化",
                Cm(23.2), Cm(9.15), Cm(8.6), Cm(0.7), font_size=16,
                color=COLORS["teal"], bold=True, align=PP_ALIGN.CENTER)
    add_page_no(slide, 2)

    # Slide 3
    slide = blank_slide(prs)
    add_title(slide, "数据与评价：同时看插值与新药泛化", "输入为药物A分子图、药物B分子图和细胞系表达；输出包括协同分类与Loewe回归")
    add_metric(slide, "160,228", "有效样本", Cm(1.35), Cm(3.25), Cm(5.3), COLORS["teal"])
    add_metric(slide, "1,896", "唯一药物", Cm(7.1), Cm(3.25), Cm(5.3), COLORS["navy"])
    add_metric(slide, "156", "细胞系", Cm(12.85), Cm(3.25), Cm(5.3), COLORS["orange"])
    add_metric(slide, "14.8%", "协同正例比例", Cm(18.6), Cm(3.25), Cm(5.3), COLORS["red"])
    add_card(slide, Cm(1.35), Cm(6.65), Cm(14.7), Cm(8.4), "两类任务", COLORS["teal"])
    add_multiline(slide, [
        "分类：预测协同/非协同，重点报告AUPR、AUC",
        "回归：预测连续Loewe score，报告RMSE、MAE",
        "AUPR更适合正例稀少的候选筛选场景"
    ], Cm(1.85), Cm(7.85), Cm(13.6), Cm(4.6), font_size=17, bullet=True)
    add_card(slide, Cm(17.25), Cm(6.65), Cm(15.2), Cm(8.4), "两种划分", COLORS["navy"])
    add_multiline(slide, [
        "random split：评估已有药物空间内的插值拟合",
        "cold-drug split：测试药物不出现在训练集中",
        "严格cold-drug：train-test药物重叠数为0"
    ], Cm(17.75), Cm(7.85), Cm(14.0), Cm(4.6), font_size=17, bullet=True)
    add_page_no(slide, 3)

    # Slide 4
    slide = blank_slide(prs)
    add_title(slide, "方法框架：比较KAN的不同放置位置", "同一输入、同一任务下，对比直接KAN head、DrugKAN、三元超图与药物对聚合图")
    add_picture_fit(slide, FIG_DIR / "model_architecture_overview.png",
                    Cm(1.35), Cm(3.0), Cm(18.2), Cm(10.0))
    add_card(slide, Cm(20.25), Cm(3.0), Cm(12.2), Cm(10.0), "模型变体", COLORS["teal"])
    add_multiline(slide, [
        "MLP baseline：药物A、药物B、细胞向量拼接后预测",
        "KAN head：直接用KAN替换最终MLP",
        "DrugKAN：在药物分子图消息传递中使用KAN",
        "HgKAN-HG：drug-drug-cell三元超图交互",
        "HgKAN-Agg：先构造drug-pair节点，再与cell交互"
    ], Cm(20.75), Cm(4.1), Cm(11.2), Cm(6.7), font_size=14.5, bullet=True)
    add_textbox(slide, "最终主结果：HgKAN-Agg",
                Cm(20.8), Cm(11.45), Cm(11.0), Cm(0.7), font_size=18,
                color=COLORS["teal"], bold=True, align=PP_ALIGN.CENTER)
    add_page_no(slide, 4)

    # Slide 5
    slide = blank_slide(prs)
    add_title(slide, "实验结果：HgKAN-Agg在cold-drug上最明确", "random split中MLP仍很强；真正有价值的是未见药物泛化下的方向一致改善")
    add_picture_fit(slide, METRIC_FIG, Cm(1.35), Cm(3.05), Cm(15.9), Cm(7.2))
    table_data = [
        ["模型", "AUPR↑", "AUC↑", "RMSE↓", "MAE↓"],
        ["MLP baseline", "0.2758", "0.7333", "18.3997", "13.6750"],
        ["KAN head", "0.1717", "0.5590", "19.2865", "14.1486"],
        ["DrugKAN", "0.1967", "0.6485", "17.8930", "13.1844"],
        ["HgKAN-Agg", "0.2765", "0.7452", "17.6763", "12.8615"],
    ]
    add_table(slide, 5, 5, Cm(18.0), Cm(3.15), Cm(14.35), Cm(5.8), table_data,
              col_widths=[Cm(4.2), Cm(2.45), Cm(2.45), Cm(2.65), Cm(2.6)],
              font_size=11)
    add_card(slide, Cm(18.0), Cm(9.45), Cm(14.35), Cm(4.3), "结论读法", COLORS["orange"])
    add_multiline(slide, [
        "直接KAN head显著退化：高维拼接表示不适合无结构KAN",
        "DrugKAN对回归有帮助，但分类协同判别不稳定",
        "HgKAN-Agg四个主指标方向一致优于MLP baseline"
    ], Cm(18.45), Cm(10.45), Cm(13.4), Cm(2.7), font_size=14.5, bullet=True)
    add_textbox(slide, "注意：提升幅度有限，因此本文强调“适用边界与结构趋势”，不夸大为全面SOTA。",
                Cm(1.55), Cm(15.1), Cm(30.3), Cm(0.7), font_size=13,
                color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_page_no(slide, 5)

    # Slide 6
    slide = blank_slide(prs)
    add_title(slide, "案例解释：模型确实依赖药物对信息", "Lenalidomide + mitomycin C @ IGROV1，cold-drug测试样本，两种药物均未出现在训练集中")
    add_picture_fit(slide, FIG_DIR / "case_57954_visual_explanation_overview.png",
                    Cm(1.35), Cm(3.0), Cm(16.0), Cm(11.0))
    add_card(slide, Cm(18.2), Cm(3.0), Cm(14.25), Cm(11.0), "关键观察", COLORS["teal"])
    add_multiline(slide, [
        "真实标签：协同；Loewe score = 12.3335",
        "模型基准协同概率：0.9106",
        "屏蔽两种药物后降至0.3711",
        "药物对表示恢复到15%时，概率升至0.6506并越过阈值",
        "原子saliency与glutarimide、quinone/mitosene等已知相关结构有重合"
    ], Cm(18.7), Cm(4.05), Cm(13.2), Cm(6.2), font_size=15, bullet=True)
    add_textbox(slide, "解释边界：这些是模型内部依赖性证据，不能直接替代真实药理机制验证。",
                Cm(18.75), Cm(11.7), Cm(12.9), Cm(1.1), font_size=13,
                color=COLORS["red"], bold=True, align=PP_ALIGN.CENTER)
    add_page_no(slide, 6)

    # Slide 7
    slide = blank_slide(prs)
    add_title(slide, "结论与展望", "本文贡献在于系统辨析KAN在药物协同预测中的有效放置方式")
    add_card(slide, Cm(1.35), Cm(3.1), Cm(9.8), Cm(10.3), "主要结论", COLORS["teal"])
    add_multiline(slide, [
        "KAN不是简单替换MLP即可稳定提升的通用模块",
        "单药图编码中的KAN只解决部分结构表示问题",
        "药物对聚合图让KAN在更合适的局部结构中学习非线性消息函数"
    ], Cm(1.85), Cm(4.15), Cm(8.8), Cm(6.3), font_size=16, bullet=True)
    add_card(slide, Cm(12.05), Cm(3.1), Cm(9.8), Cm(10.3), "本文贡献", COLORS["navy"])
    add_multiline(slide, [
        "构建DrugComb药物协同预测数据与严格cold-drug评估",
        "比较KAN head、DrugKAN、HgKAN-HG、HgKAN-Agg等多种变体",
        "建立KAN曲线、扰动、药物对响应和原子saliency解释流程"
    ], Cm(12.55), Cm(4.15), Cm(8.8), Cm(6.3), font_size=16, bullet=True)
    add_card(slide, Cm(22.75), Cm(3.1), Cm(9.7), Cm(10.3), "不足与后续", COLORS["orange"])
    add_multiline(slide, [
        "cold-drug分类仍存在明显划分方差",
        "多随机种子和外部数据集验证还需加强",
        "未来可加入剂量变量、多组学细胞表征和更系统的解释验证"
    ], Cm(23.25), Cm(4.15), Cm(8.7), Cm(6.3), font_size=16, bullet=True)
    add_textbox(slide, "最终观点：结构先行，KAN作为结构化消息传递中的可学习函数模块更有价值。",
                Cm(3.3), Cm(15.0), Cm(27.4), Cm(0.9), font_size=20,
                color=COLORS["teal"], bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, "谢谢各位老师，欢迎批评指正",
                Cm(10.4), Cm(16.35), Cm(13.0), Cm(0.75), font_size=18,
                color=COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    add_page_no(slide, 7)

    prs.save(PPTX_PATH)
    add_notes()


if __name__ == "__main__":
    build_deck()
    print(PPTX_PATH)
    print(NOTES_PATH)
